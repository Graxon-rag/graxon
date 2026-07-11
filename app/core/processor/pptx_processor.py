from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from app.utils.logger import logger
from .processor import Processor
from app.config.env import Env
from typing import List, Tuple
from pptx import Presentation
import os


class PPTXProcessor(Processor):
    def __init__(self,
        file_path: str,
        filename: str,
        chunk_number: int,
        rag_chunk_start_index: int,
        pages_per_batch: int = 20,          # maps directly to slides_per_batch
        rag_chunk_size: int = Env.CHUNK_SIZE,
        rag_chunk_overlap: int = Env.CHUNK_OVERLAP,
        # no tail_carry_chars — slides are independent units
    ):
        self.file_path = file_path
        self.filename = filename
        self.chunk_number = chunk_number
        self.rag_chunk_start_index = rag_chunk_start_index
        self.pages_per_batch = pages_per_batch
        self.rag_chunk_size = rag_chunk_size
        self.rag_chunk_overlap = rag_chunk_overlap

        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=rag_chunk_size,
            chunk_overlap=rag_chunk_overlap,
            separators=["\n\n", "\n", ".", " ", ""],
        )

    def _extract_slide_text(self, slide) -> str:
        """
        Extracts all text from a slide's shapes and text frames.
        Preserves title first, then body content.
        """
        texts = []
        title = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # picture, skip
                continue

            # extract title separately so it always comes first
            if hasattr(shape, "name") and "title" in shape.name.lower():
                title = shape.text_frame.text.strip()
            else:
                shape_text = "\n".join(
                    para.text.strip()
                    for para in shape.text_frame.paragraphs
                    if para.text.strip()
                )
                if shape_text:
                    texts.append(shape_text)

        result = []
        if title:
            result.append(title)
        result.extend(texts)

        return "\n".join(result)

    async def process(self) -> Tuple[List[Document], int, bool]:
        """
        Reads pages_per_batch slides starting from chunk_number * pages_per_batch.
        No tail carry — slides are independent units.

        Returns:
            documents: list of Document
            next_rag_chunk_start_index: pass this to next queue message
            is_last: True if this was the final batch
        """
        try:
            prs = Presentation(self.file_path)
            total_slides = len(prs.slides)

            slide_start = self.chunk_number * self.pages_per_batch
            if slide_start >= total_slides:
                raise ValueError(
                    f"chunk_number {self.chunk_number} is out of range. "
                    f"Total slides: {total_slides}, pages_per_batch: {self.pages_per_batch}"
                )

            slide_end = min(slide_start + self.pages_per_batch, total_slides)
            is_last = slide_end >= total_slides

            # extract text from each slide in batch
            batch_text = ""
            for slide_num in range(slide_start, slide_end):
                slide_text = self._extract_slide_text(prs.slides[slide_num])
                if slide_text:
                    batch_text += slide_text + "\n\n"   # double newline between slides

            documents = self._split_into_rag_chunks(batch_text.strip(), slide_start)
            return documents, self.rag_chunk_start_index + len(documents), is_last

        except Exception as e:
            logger.error(f"Failed to process PPTX file {self.file_path}. Error: {e}")
            raise e

    def _split_into_rag_chunks(self, raw_text: str, slide_start: int) -> List[Document]:
        texts = self.splitter.split_text(raw_text)

        documents = []
        for i, text in enumerate(texts):
            absolute_index = self.rag_chunk_start_index + i
            doc = Document(
                id=f"{self.filename}-{absolute_index}",
                page_content=text,
                metadata={
                    "source": self.file_path,
                    "file_chunk_number": self.chunk_number,
                    "rag_chunk_number": absolute_index,
                    "slide_start": slide_start,         # which slide batch this came from
                },
            )
            documents.append(doc)

        return documents
